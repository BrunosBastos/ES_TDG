from fastapi import FastAPI, UploadFile, File, Depends, Form, Header
from fastapi.responses import JSONResponse
from pptx import Presentation
from botocore.exceptions import ClientError
from fastapi.middleware.cors import CORSMiddleware
import logging
import openpyxl as oxl
import json
import re
from utils import create_response,    \
                  delete_paragraph,   \
                  get_client_s3,      \
                  get_file_extension, \
                  duplicate_slide,    \
                  fill_paragraph,     \
                  insert_paragraph
import docx
from typing import Optional

app = FastAPI()
bucket_name = "tdg-s3-bucket"

# Logger
logging.basicConfig(
    format="%(module)-20s:%(levelname)-15s| %(message)s",
    level=logging.INFO
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.post("/api/3/fill")
async def post_template(
        upload_file: UploadFile = File(...),
        retrieval_filename: str = Form(),
        output_filename: str = Form(),
        username: Optional[str] = Header(None),
        s3=Depends(get_client_s3)
) -> JSONResponse:
    """
    Endpoint ``/fill`` that accepts the method POST. Receives a JSON file
    and a template name.

    Parameters
    ----------
        upload_file : `UploadFile`
            The file provided in the POST request
        retrieval_filename : `str`
            Location of the template file in the bucket
        output_filename : `str`
            The name of the resulting file

    Returns
    -------
        response : `JSONResponse`
            Json response with the status code and data containing the message
            and data.

    """
    try:

        if not username:
            return create_response(status_code=401, message="Not authenticated.")

        _, t_format, t_filename = retrieval_filename.split("/")

        retrieval_filename = username + "/" + retrieval_filename

        try:
            # load json data
            data = json.load(upload_file.file)
        except Exception:
            return create_response(status_code=400,
                                   message="Invalid Json file format.")

        s3.download_file(
            bucket_name, retrieval_filename, t_filename
        )

        file_extension = get_file_extension(t_filename, t_format)

        if not output_filename.endswith("." + file_extension):
            output_filename += "." + file_extension

        if not fill_template(file_extension, t_filename, data, output_filename):
            return create_response(status_code=400,
                                   message="Could not fill template it the provided data." +
                                   "Check if the file matches the requirements.")

        path = username + "/filled/" + t_format + "/" + output_filename

        s3.upload_fileobj(
            open(output_filename, "rb"), bucket_name, path
        )

        return create_response(status_code=200, data=output_filename)

    except ClientError as e:
        logging.debug(e)
        return create_response(status_code=400, message=e)


def fill_template(file_extension, t_filename, data, output_filename):
    """
    Chooses the right function to handle the template filling task.

    Parameters
    ----------
        file_extension : `str`
            The extension of the file
        t_filename : `str`
            Name of the template
        data : `json`
            Json data already converted to a dict
        output_filename : `str`
            Name of the resulting file

    Returns
    -------
        valid : `bool`
            True if there is no error, else False
    """

    # fill excel template
    if file_extension == "xlsx":
        return fill_excel_template(t_filename, data, output_filename)

    # fill pptx template
    elif file_extension == "pptx" or file_extension == "ppt":
        return fill_ppt_template(t_filename, data, output_filename)

    # fill pptx template
    elif file_extension == "docx":
        return fill_docx_template(t_filename, data, output_filename)


def fill_excel_template(template_name, data, filled_file_name):
    """
    Fills excel template with JSON data

    Parameters
    ----------
        template_name : `str`
            The name of the template file
        data : JSON
            The JSON data to be used in filling
        filled_file_name : `str`
            The name of the resulting file
    Returns
    -------
        valid: `bool`
            True if there was no error else False
    """

    try:
        # load template
        template = oxl.load_workbook(template_name)

        i = 0
        # go through all sheets
        for sheet in template.sheetnames:

            if i > 0 and sheet not in template:
                # create new sheet
                ws = template.create_sheet(index=i, title=sheet)
            else:
                ws = template[sheet]

            # fill sheet with data
            for cell in data[i][sheet]:
                cell_number = list(cell.keys())
                ws[cell_number[0]] = cell[cell_number[0]]

            i += 1

        template.save(filled_file_name)
        return True
    except Exception:
        return False


def fill_docx_template(template_name, data, filled_file_name):  # noqa: C901
    """
    Fills word template with JSON data

    Parameters
    ----------
        template_name : `str`
            The name of the template file
        data : JSON
            The JSON data to be used in filling
        filled_file_name : `str`
            The name of the resulting file
    Returns
    -------
        valid: `bool`
            True if there was no error else False
    """

    try:
        # open file
        document = docx.Document(template_name)

        # regex
        begin_list_regex = re.compile(r"\$\{#\w+}*")
        end_list_regex = re.compile(r"\$\{\w+#}*")
        list_value_regex = re.compile(r"\$\{\.\w+}*")

        # iterate through paragraphs in document
        in_list = False
        in_object_name = None

        # stores the paragraphs that are inside a list
        lst_paragraphs = []
        del_paragraphs = []
        idx_counter = 0

        for pi, paragraph in enumerate(document.paragraphs):

            local_data = None
            if in_list:
                local_data = data[in_object_name][0]
                lst_paragraphs.append(paragraph.text)

            paragraph.text = fill_paragraph(data, local_data, paragraph.text)

            # find's the beginning of a list
            if x := re.search(begin_list_regex, paragraph.text):
                in_object_name = x.string[3:-1]
                print(in_object_name)
                in_list = True
                del_paragraphs.append(paragraph)

            # find's the end of a list
            if x := re.search(end_list_regex, paragraph.text):
                print(in_object_name)

                del_paragraphs.append(paragraph)
                lst_paragraphs.pop()        # remove the last value because its the closing tag

                for value in data[in_object_name][1:]:
                    for p in lst_paragraphs:
                        insert_paragraph(
                            document,
                            pi + idx_counter,
                            fill_paragraph(data, value, p),
                        )
                        idx_counter += 1
                in_object_name = None
                in_list = False
                lst_paragraphs.clear()

        # deletes the unused paragraphs
        for p in del_paragraphs:
            delete_paragraph(p)

        del_paragraphs = []

        # iterate through tables in document
        # find json table names to iterate through
        table_names = []
        for idx, table in enumerate(document.tables):
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if x := re.search(begin_list_regex, paragraph.text):
                            table_names.append(x.group(0))

        # add enough rows and cells to the table according to json size
        cell_info = []
        for idx, table in enumerate(document.tables):
            name = table_names[idx][3:-1]

            # add missing rows
            while len(table.rows) < len(data[name]) + 1:
                table.add_row()

            # get cell info for replication
            for cell in table.rows[1].cells:
                text = [
                    paragraph.text for paragraph in cell.paragraphs if paragraph.text != ""]
                cell_info.append(text)

            # replicate info on each cell
            for row in table.rows[1:]:
                for x, cell in enumerate(row.cells):
                    cell.paragraphs[0].text = cell_info[x]

            # fill table rows with json data
            for x, row in enumerate(table.rows[1:]):
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if replace := re.search(begin_list_regex, paragraph.text):
                            paragraph.text = paragraph.text.replace(
                                replace.group(0), "")

                        if replace := re.search(end_list_regex, paragraph.text):
                            paragraph.text = paragraph.text.replace(
                                replace.group(0), "")

                        if replace := re.search(list_value_regex, paragraph.text):
                            paragraph.text = paragraph.text.replace(
                                replace.group(0),
                                data[table_names[idx][3:-1]
                                     ][x][replace.group(0)[3:-1]]
                            )

        document.save(filled_file_name)
        return True
    except Exception:
        return False


def fill_ppt_template(template_name, data, filled_file_name):
    """
    Fills powerpoint template with JSON data

    Parameters
    ----------
        template_name : `str`
            The name of the template file
        data : JSON
            The JSON data to be used in filling
        filled_file_name : `str`
            The name of the resulting file
    Returns
    -------
        valid: `bool`
            True if there was no error else False
    """

    try:
        # load template
        template = Presentation(template_name)

        [duplicate_slide(template, 0) for _ in range(1, len(data["records"]))]

        i = 0
        # go through data
        for slide_data in data["records"]:
            slide_shape = template.slides[i]
            for slide_t in slide_shape.shapes:
                slide_text = slide_t.text_frame
                for slide_paragraph in slide_text.paragraphs:
                    whole_text = "".join(
                        [r.text for r in slide_paragraph.runs])
                    data_params = re.findall(
                        r"\{([A-Za-z0-9_]+)\}", whole_text)
                    for d in data_params:
                        if d in slide_data.keys():
                            whole_text = whole_text.replace(
                                f"{{{d}}}", slide_data[d])
                    for idx, run in enumerate(slide_paragraph.runs):
                        if idx == 0:
                            continue
                        p = slide_paragraph._p
                        p.remove(run._r)
                    slide_paragraph.runs[0].text = whole_text

            i += 1

        template.save(filled_file_name)
        return True
    except Exception:
        return False
